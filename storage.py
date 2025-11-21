import os
import json
import time
import fitz  # PyMuPDF (필요시 유지, 여기서는 fallback 용도)
import pptx
from flask import session, current_app
import google.generativeai as genai # [!! 중요 !!] Gemini 사용

# [!! 중요 !!] app.py에서 data_lock을 가져오되, 순환 참조 방지
try:
    from app import data_lock
except ImportError:
    import threading
    data_lock = threading.RLock()

# 설정값
BASE_DATA_DIR = "data"
BASE_CACHE_DIR = "cache"
ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'png', 'jpg', 'jpeg', 'txt', 'xlsx'}

def get_user_data_path(user_id):
    """사용자의 data 폴더 경로를 반환합니다."""
    path = os.path.join(BASE_DATA_DIR, user_id)
    os.makedirs(path, exist_ok=True)
    return path

def get_user_cache_path(user_id, cache_type="qa"):
    """사용자의 cache 파일 경로를 반환합니다."""
    os.makedirs(BASE_CACHE_DIR, exist_ok=True)
    return os.path.join(BASE_CACHE_DIR, f"{cache_type}_{user_id}.json")

def load_qa_cache(user_id):
    """Q&A 캐시 로드"""
    qa_cache_file = get_user_cache_path(user_id, "qa")
    if os.path.exists(qa_cache_file):
        try:
            with open(qa_cache_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"⚠️ [Cache] '{user_id}' Q&A 로드 실패: {e}")
            return {}
    return {}

def save_qa_cache(user_id, qa_cache):
    """Q&A 캐시 저장"""
    qa_cache_file = get_user_cache_path(user_id, "qa")
    with data_lock:
        try:
            with open(qa_cache_file, 'w', encoding='utf-8') as f:
                json.dump(qa_cache, f, ensure_ascii=False, indent=4)
        except Exception as e:
            print(f"💥 [Cache] '{user_id}' Q&A 저장 실패: {e}")

def load_ocr_cache(user_id):
    """OCR 캐시 로드"""
    ocr_cache_file = get_user_cache_path(user_id, "ocr")
    if os.path.exists(ocr_cache_file):
        try:
            with open(ocr_cache_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"⚠️ [Cache] '{user_id}' OCR 로드 실패: {e}")
            return {}
    return {}

def save_ocr_cache(user_id, ocr_cache):
    """OCR 캐시 저장"""
    ocr_cache_file = get_user_cache_path(user_id, "ocr")
    with data_lock:
        try:
            with open(ocr_cache_file, 'w', encoding='utf-8') as f:
                json.dump(ocr_cache, f, ensure_ascii=False, indent=4)
        except Exception as e:
            print(f"💥 [Cache] '{user_id}' OCR 저장 실패: {e}")

def load_odapnote(user_id):
    """오답노트 로드"""
    odapnote_file = get_user_cache_path(user_id, "odap")
    if os.path.exists(odapnote_file):
        try:
            with open(odapnote_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"⚠️ [Cache] '{user_id}' 오답노트 로드 실패: {e}")
            return []
    return []

def save_odapnote(user_id, odapnote_list):
    """오답노트 저장"""
    odapnote_file = get_user_cache_path(user_id, "odap")
    with data_lock:
        try:
            with open(odapnote_file, 'w', encoding='utf-8') as f:
                json.dump(odapnote_list, f, ensure_ascii=False, indent=4)
        except Exception as e:
            print(f"💥 [Cache] '{user_id}' 오답노트 저장 실패: {e}")

def get_supported_files(user_id):
    """사용자 폴더의 지원되는 파일 목록 반환"""
    user_data_path = get_user_data_path(user_id)
    if not os.path.exists(user_data_path):
        return []
    return sorted([f for f in os.listdir(user_data_path) 
                   if any(f.lower().endswith(ext) for ext in ALLOWED_EXTENSIONS)])

def allowed_file(filename):
    """확장자 체크"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_text_from_single_file(user_id, filename, force_ocr=False):
    """
    [!! 핵심 변경 !!] 
    Gemini Vision API를 사용하여 PDF/이미지 텍스트를 추출합니다.
    서버에 Tesseract를 설치할 필요가 없습니다.
    """
    user_data_path = get_user_data_path(user_id)
    ocr_cache = load_ocr_cache(user_id)
    file_path = os.path.join(user_data_path, filename)
    
    # 1. 캐시 확인
    if not force_ocr:
        with data_lock:
            cached_text = ocr_cache.get(filename)
        if cached_text is not None:
            print(f"⚡️ [OCR 캐시 HIT] '{user_id}/{filename}' 로드 완료.")
            return cached_text
            
    print(f"🧠 [Gemini OCR] '{user_id}/{filename}' 처리 시작 (Google 서버로 전송)...")

    if not os.path.exists(file_path):
        return None
        
    full_text = ""
    try:
        # ==================================================
        # (A) PDF 또는 이미지 파일 -> Gemini에게 통째로 맡김
        # ==================================================
        if filename.lower().endswith(('.pdf', '.png', '.jpg', '.jpeg')):
            
            # 1. 파일 업로드 (Gemini 서버로)
            print(f"   - [1/3] '{filename}' 업로드 중...")
            sample_file = genai.upload_file(path=file_path, display_name=filename)
            
            # 2. 파일 처리 대기 (Active 상태 될 때까지 폴링)
            print(f"   - [2/3] 구글 서버 처리 대기 중...")
            while sample_file.state.name == "PROCESSING":
                time.sleep(1)
                sample_file = genai.get_file(sample_file.name)

            if sample_file.state.name == "FAILED":
                raise ValueError("Gemini 서버에서 파일 처리에 실패했습니다.")

            # 3. 텍스트 추출 요청 (Flash 모델 사용)
            print("   - [3/3] 텍스트 추출 요청 중...")
            model = genai.GenerativeModel("gemini-1.5-flash") 
            response = model.generate_content([
                "이 파일에 있는 모든 텍스트를 처음부터 끝까지 순서대로, 빠짐없이 추출해줘. 요약하지 말고 원문 텍스트만 그대로 줘. 불필요한 설명은 생략해.", 
                sample_file
            ])
            
            full_text = response.text
            
            # 4. (중요) Gemini 서버에서 파일 삭제 (공간/비용 절약)
            try:
                genai.delete_file(sample_file.name)
                print("   - (청소) Gemini 서버 임시 파일 삭제 완료.")
            except:
                pass

        # ==================================================
        # (B) PPTX, TXT, Excel -> 로컬 파이썬 라이브러리 사용 (빠르고 무료)
        # ==================================================
        elif filename.lower().endswith('.pptx'):
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
        
        elif filename.lower().endswith('.txt'):
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    full_text = f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='cp949') as f:
                    full_text = f.read()
        
        elif filename.lower().endswith('.xlsx'):
            if not current_app.config.get('OPENPYXL_AVAILABLE', False):
                print("⚠️ openpyxl 라이브러리가 없어 엑셀 파일을 읽을 수 없습니다.")
                return None 
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row in sheet.iter_rows():
                    row_text = [str(cell.value) for cell in row if cell.value is not None]
                    full_text += " ".join(row_text) + "\n" 
            wb.close()

        # 결과 저장 및 반환
        print(f"✅ [Gemini OCR] '{user_id}/{filename}' 완료 (길이: {len(full_text)}자).")
        
        with data_lock:
            ocr_cache[filename] = full_text
            save_ocr_cache(user_id, ocr_cache) 
        
        return full_text

    except Exception as e:
        print(f"❌ [Gemini OCR 오류] '{user_id}/{filename}': {e}")
        # 오류 시 캐시 저장 안 함 (재시도 가능하게)
        return None

def load_all_text_from_data(user_id):
    """OCR 캐시에서 사용자의 모든 텍스트 로드"""
    print(f"🔄 [Storage] '{user_id}' 통합 텍스트 구성 중...")
    temp_text_list = []
    
    current_files = get_supported_files(user_id)
    ocr_cache = load_ocr_cache(user_id)
    
    with data_lock:
        cache_updated = False
        cached_files = list(ocr_cache.keys())
        # 파일 삭제 동기화
        for f in cached_files:
            if f not in current_files:
                del ocr_cache[f]
                cache_updated = True
        
        if cache_updated:
            save_ocr_cache(user_id, ocr_cache)

        for filename in current_files:
            text = ocr_cache.get(filename) 
            if text: 
                temp_text_list.append(f"--- {filename} 시작 ---\n{text}\n--- {filename} 끝 ---")
                
    all_file_text = "\n\n".join(temp_text_list)
    return all_file_text

def get_categorized_cache(qa_cache):
    """Q&A 캐시 분류"""
    ask_list = []
    summarize_list = []
    quiz_list = []
    mindmap_list = []

    sorted_items = sorted(qa_cache.items(), key=lambda item: item[1].get('timestamp', '0'), reverse=True)

    for key, value in sorted_items:
        action_type = value.get('action_type', 'ask')
        if action_type in ['ask', 'quiz_file']:
            ask_list.append({'key': key, 'value': value})
        elif action_type in ['extract_answer', 'extract_all']:
            summarize_list.append({'key': key, 'value': value})
        elif action_type in ['quiz_all', 'quiz_selected', 'quiz_weakness', 'grade_quiz', 'analyze_weakness']:
            quiz_list.append({'key': key, 'value': value})
        elif action_type == 'generate_mindmap':
            mindmap_list.append({'key': key, 'value': value})

    return ask_list, summarize_list, quiz_list, mindmap_list